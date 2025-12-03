import io
import math
import streamlit as st
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
import matplotlib.pyplot as plt
import matplotlib.patches as patches

# ============================================================
# 1. UTILS & CONFIGURATION
# ============================================================

def get_tx_chain(voltage, scheme):
    if voltage == "400V": return []
    if voltage == "11kV": return [{"ratio": "11/0.4 kV", "bus": "0.4kV"}]
    if voltage == "33kV": return [{"ratio": "33/0.4 kV", "bus": "0.4kV"}]
    if voltage == "132kV":
        if scheme == "132/0.4 kV": return [{"ratio": "132/0.4 kV", "bus": "0.4kV"}]
        elif scheme == "132/11/0.4 kV": return [{"ratio": "132/11 kV", "bus": "11kV"}, {"ratio": "11/0.4 kV", "bus": "0.4kV"}]
        else: return [{"ratio": "132/33 kV", "bus": "33kV"}, {"ratio": "33/11 kV", "bus": "11kV"}, {"ratio": "11/0.4 kV", "bus": "0.4kV"}]
    return []

def get_lv_gen_inputs(key_prefix, include_emsb=False):
    if include_emsb:
        c1, c2, c3 = st.columns(3)
    else:
        c1, c2 = st.columns(2)
        
    has_solar = c1.checkbox("Solar PV", key=f"{key_prefix}_sol")
    has_bess = c2.checkbox("BESS", key=f"{key_prefix}_bess")
    
    has_emsb = False
    if include_emsb:
        has_emsb = c3.checkbox("EMSB", key=f"{key_prefix}_emsb")
    
    gens = []
    if has_solar:
        st.markdown("**Solar PV Specs**")
        ca, cb = st.columns(2)
        kwac = ca.number_input("kWac", 0, 99999, 100, key=f"{key_prefix}_skwa")
        kwp = cb.number_input("kWp", 0, 99999, 120, key=f"{key_prefix}_skwp")
        gens.append({"type": "Solar", "kWac": kwac, "cap_val": kwp})
        
    if has_bess:
        st.markdown("**BESS Specs**")
        ca, cb = st.columns(2)
        kwac = ca.number_input("kWac", 0, 99999, 100, key=f"{key_prefix}_bkwa")
        kwh = cb.number_input("kWh", 0, 99999, 200, key=f"{key_prefix}_bkwh")
        gens.append({"type": "BESS", "kWac": kwac, "cap_val": kwh})
        
    return gens, has_emsb

def get_mv_gen_inputs(key_prefix):
    st.markdown("**MV Generation Source**")
    gen_type = st.radio("Source Type", ["Solar PV", "BESS"], horizontal=True, key=f"{key_prefix}_type")
    
    gens = []
    c1, c2 = st.columns(2)
    
    if gen_type == "Solar PV":
        kwac = c1.number_input("kWac", 0, 99999, 1000, key=f"{key_prefix}_mv_skwa")
        kwp = c2.number_input("kWp", 0, 99999, 1200, key=f"{key_prefix}_mv_skwp")
        gens.append({"type": "Solar", "kWac": kwac, "cap_val": kwp})
    else:
        kwac = c1.number_input("kWac", 0, 99999, 1000, key=f"{key_prefix}_mv_bkwa")
        kwh = c2.number_input("kWh", 0, 99999, 2000, key=f"{key_prefix}_mv_bkwh")
        gens.append({"type": "BESS", "kWac": kwac, "cap_val": kwh})
        
    return gens

def get_feeder_width_config(is_pptx=True):
    return {
        "item_w": 3.5,  
        "min_w": 5.0,   
        "gap": 1.0,     
        "sub_gap": 1.0  
    }

def calculate_single_feeder_width(config, dims):
    c_type = config.get("type", "Standard")
    
    if c_type == "Sub-Board":
        sub_feeders = config.get("sub_feeders", {})
        n_subs = len(sub_feeders)
        if n_subs == 0: 
            return dims["min_w"], []
            
        sub_widths = []
        for j in range(n_subs):
            s_conf = sub_feeders.get(j, {})
            sf_type = s_conf.get("type", "Standard")
            
            if sf_type == "MV Gen":
                calc_w = dims["item_w"]
            elif sf_type == "Extension":
                ext_feeders = s_conf.get("extension_feeders", {})
                n_ext = len(ext_feeders)
                if n_ext == 0:
                    calc_w = dims["item_w"]
                else:
                    total_ext_w = 0
                    for k in ext_feeders:
                        ef_conf = ext_feeders[k]
                        ef_type = ef_conf.get("type", "Standard")
                        if ef_type == "MV Gen":
                            ef_w = dims["item_w"]
                        else:
                            s_gens = ef_conf.get("gens", [])
                            s_emsb = ef_conf.get("has_emsb", False)
                            item_count = len(s_gens) + (1 if s_emsb else 0)
                            ef_w = max(dims["item_w"] * 1.5, item_count * dims["item_w"])
                        total_ext_w += ef_w
                    total_ext_w += (n_ext - 1) * (dims["sub_gap"] * 0.8)
                    calc_w = max(dims["item_w"], total_ext_w)
            else:
                s_gens = s_conf.get("gens", [])
                s_emsb = s_conf.get("has_emsb", False)
                item_count = len(s_gens) + (1 if s_emsb else 0)
                calc_w = max(dims["item_w"] * 1.5, item_count * dims["item_w"])
            
            sub_widths.append(calc_w)
            
        total_sub_width = sum(sub_widths) + (len(sub_widths) - 1) * dims["sub_gap"]
        final_total_w = max(dims["min_w"], total_sub_width)
        return final_total_w, sub_widths
        
    else:
        gens = config.get("gens", [])
        has_emsb = config.get("emsb", {}).get("has", False)
        item_count = len(gens) + (1 if has_emsb else 0)
        
        if item_count <= 1:
            return dims["min_w"], []
        else:
            return max(dims["min_w"], item_count * dims["item_w"]), []

def calculate_section_layout(section_feeders, swg_configs, start_x, is_pptx=False):
    dims = get_feeder_width_config(is_pptx)
    current_x = start_x
    feeder_centers = []
    feeder_widths = []
    sub_widths_map = {}
    
    if not section_feeders: return 0, [], [], {}

    for i in section_feeders:
        config = swg_configs.get(i, {})
        w, sub_ws = calculate_single_feeder_width(config, dims)
        
        center = current_x + (w / 2)
        feeder_centers.append(center)
        feeder_widths.append(w)
        sub_widths_map[i] = sub_ws
        
        current_x += w + dims["gap"]
        
    total_width = current_x - start_x - dims["gap"] 
    return total_width, feeder_centers, feeder_widths, sub_widths_map

# ============================================================
# 2. PPTX DRAWING HELPERS
# ============================================================

def S(val, scale):
    """Safe scaling to Inches, returns float-like Inch object"""
    return Inches(float(val) * float(scale))

def add_line(slide, x1, y1, x2, y2, width_pt=3, color=RGBColor(0, 112, 192)):
    # Safely cast to int (EMU) for PPTX
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    conn.line.width = Pt(width_pt)
    conn.line.color.rgb = color
    return conn

def add_busbar(slide, left, top, width):
    if width <= 0: return
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, int(left), int(top), int(width), Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 112, 192)
    bar.line.fill.background()

def add_breaker_x(slide, cx, cy, scale, size_base=0.25, color=RGBColor(0, 112, 192)):
    # cx, cy are expected to be int (EMU)
    half = int(S(size_base, scale))
    add_line(slide, cx - half, cy - half, cx + half, cy + half, 3.0, color)
    add_line(slide, cx - half, cy + half, cx + half, cy - half, 3.0, color)

def pptx_add_transformer(slide, cx_int, center_y, ratio_txt, tx_id, scale):
    # cx_int, center_y expected as int (EMU)
    r = int(S(0.35, scale)); d = r * 2
    top_y = center_y - r; bot_y = center_y + r
    
    s1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, int(cx_int - r), int(top_y - r), int(d), int(d))
    s1.fill.solid(); s1.fill.fore_color.rgb = RGBColor(255, 255, 255) 
    s1.line.width = Pt(2.0); s1.line.color.rgb = RGBColor(0,0,0)
    
    s2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, int(cx_int - r), int(bot_y - r), int(d), int(d))
    s2.fill.solid(); s2.fill.fore_color.rgb = RGBColor(255, 255, 255) 
    s2.line.width = Pt(2.0); s2.line.color.rgb = RGBColor(0,0,0)
    
    line_y1 = int(center_y - S(0.9, scale))
    line_y2 = int(top_y - r + S(0.05, scale))
    line_y3 = int(bot_y + r - S(0.05, scale))
    line_y4 = int(center_y + S(0.9, scale))
    
    add_line(slide, cx_int, line_y1, cx_int, line_y2, 3, RGBColor(0,0,0))
    add_line(slide, cx_int, line_y3, cx_int, line_y4, 3, RGBColor(0,0,0))

    tb = slide.shapes.add_textbox(int(cx_int + S(0.4, scale)), int(center_y - S(0.8, scale)), int(S(4.0, scale)), int(S(1.5, scale)))
    p = tb.text_frame.paragraphs[0]; p.text = f"{tx_id}\n{ratio_txt}"; 
    p.font.bold = True; p.font.size = Pt(max(10, 20*scale)) 

def pptx_add_inverter_branch(slide, cx, start_y, gens, scale):
    # cx, start_y as int (EMU)
    if not gens: return
    box_top = int(start_y + S(1.5, scale))
    add_line(slide, cx, start_y, cx, box_top, 2, RGBColor(0, 176, 80))
    gen = gens[0]
    w = int(S(2.2, scale)); h = int(S(1.5, scale)); left = int(cx - w/2)
    
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, box_top, w, h)
    s.fill.background(); s.line.color.rgb = RGBColor(0, 176, 80); s.line.width = Pt(2.0)
    add_line(slide, left, box_top, left + w, box_top+h, 2, RGBColor(0, 176, 80))
    add_line(slide, left, box_top+h, left + w, box_top, 2, RGBColor(0, 176, 80))
    
    title = "BESS" if gen['type'] == "BESS" else "SOLAR PV"
    cap_unit = "kWh" if gen['type'] == "BESS" else "kWp"
    lines = [title, f"{gen['kWac']} kWac", f"{gen['cap_val']} {cap_unit}"]
    
    tb = slide.shapes.add_textbox(int(cx - S(3, scale)), int(box_top + h + S(0.1, scale)), int(S(6, scale)), int(S(2.5, scale)))
    tf = tb.text_frame
    for l in lines:
        p = tf.add_paragraph(); p.text = l 
        p.font.size = Pt(max(10, 20*scale)); 
        p.font.color.rgb = RGBColor(0, 176, 80); 
        p.alignment = PP_ALIGN.CENTER; p.font.bold = True

def pptx_add_lv_system(slide, cx, start_y, gens, has_emsb, emsb_name, scale):
    # cx, start_y as int (EMU)
    if not gens and not has_emsb: return
    items = []
    for g in gens: items.append(('GEN', g))
    if has_emsb: items.append(('EMSB', emsb_name))
    num_items = len(items)
    if num_items == 0: return
    
    spacing = int(S(4.0, scale))
    total_width = (num_items - 1) * spacing
    start_x_offset = int(cx - total_width / 2)
    
    for idx, (itype, data) in enumerate(items):
        px = start_x_offset + idx * spacing
        box_top = int(start_y + S(1.5, scale))
        
        add_line(slide, px, start_y, px, int(box_top + S(0.05, scale)), 2, RGBColor(0, 176, 80))
        
        if itype == 'GEN':
            w = int(S(2.2, scale)); h = int(S(1.5, scale)); left = int(px - w/2)
            s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, box_top, w, h)
            s.fill.background(); s.line.color.rgb = RGBColor(0, 176, 80); s.line.width = Pt(2.0)
            add_line(slide, left, int(box_top + h), left + w, box_top, 2, RGBColor(0, 176, 80))
            
            title = "BESS" if data['type'] == "BESS" else "SOLAR PV"
            lines = [title, f"{data['kWac']} kWac", f"{data['cap_val']} {'kWh' if data['type']=='BESS' else 'kWp'}"]
            
            tb = slide.shapes.add_textbox(int(px - S(2.5, scale)), int(box_top + h + S(0.1, scale)), int(S(5.0, scale)), int(S(2.0, scale)))
            tf = tb.text_frame
            for l in lines:
                p = tf.add_paragraph(); p.text = l; 
                p.font.size = Pt(max(10, 20*scale)); 
                p.font.color.rgb = RGBColor(0, 176, 80); 
                p.alignment = PP_ALIGN.CENTER; p.font.bold = True
                
        elif itype == 'EMSB':
            breaker_y = int(start_y + S(0.8, scale))
            add_breaker_x(slide, px, breaker_y, scale, 0.20)
            w = int(S(1.6, scale)); h = int(S(0.8, scale)); left = int(px - w/2)
            s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, box_top, w, h)
            s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0, 112, 192); s.line.fill.background()
            
            tb = slide.shapes.add_textbox(int(px - S(1.5, scale)), int(box_top + h), int(S(3, scale)), int(S(0.8, scale)))
            tb.text_frame.text = data
            tb.text_frame.paragraphs[0].font.size = Pt(max(10, 20*scale)); 
            tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 

def add_continuation_arrow(slide, x, y, direction, label, scale):
    # x, y as int (EMU)
    w_arrow = int(S(0.5, scale))
    h_arrow = int(S(0.2, scale))
    y_pos = int(y - S(0.1, scale))
    
    if direction == "next":
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y_pos, w_arrow, h_arrow)
        text_x = int(x - S(1.5, scale))
        align = PP_ALIGN.RIGHT
    else:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, int(x - w_arrow), y_pos, w_arrow, h_arrow)
        shape.rotation = 180
        text_x = int(x + w_arrow)
        align = PP_ALIGN.LEFT

    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 0, 0); shape.line.fill.background()

    if label:
        tb = slide.shapes.add_textbox(text_x, int(y - S(1.0, scale)), int(S(2.5, scale)), int(S(0.8, scale)))
        p = tb.text_frame.paragraphs[0]; p.text = label
        p.font.size = Pt(max(10, 20*scale)); 
        p.font.bold = True; p.font.color.rgb = RGBColor(255, 0, 0); p.alignment = align

# ============================================================
# 3. MATPLOTLIB PREVIEW HELPERS
# ============================================================

def draw_tx_mpl(ax, x, y, label, ratio):
    r = 0.3
    c1 = plt.Circle((x, y + 0.25), r, fill=True, fc='white', ec='black', lw=2, zorder=10)
    c2 = plt.Circle((x, y - 0.25), r, fill=True, fc='white', ec='black', lw=2, zorder=10)
    ax.add_patch(c1); ax.add_patch(c2)
    ax.text(x + 0.35, y, f"{label}\n{ratio}", va="center", fontsize=8, fontweight='bold')

def draw_lv_system_mpl(ax, cx, start_y, gens, has_emsb, emsb_name):
    if not gens and not has_emsb: return
    items = []
    for g in gens: items.append(('GEN', g))
    if has_emsb: items.append(('EMSB', emsb_name))
    num = len(items)
    
    spacing = 3.5 
    total_w = (num - 1) * spacing
    start_x = cx - total_w / 2
    
    for idx, (itype, data) in enumerate(items):
        px = start_x + idx * spacing
        box_top = start_y - 1.5
        
        ax.plot([px, px], [start_y, box_top], color="tab:green", lw=3)
        
        if itype == 'GEN':
            rect = patches.Rectangle((px - 0.75, box_top - 1.0), 1.5, 1.0, fill=False, edgecolor="tab:green", lw=2)
            ax.add_patch(rect)
            ax.plot([px - 0.75, px + 0.75], [box_top - 1.0, box_top], color="tab:green", lw=1.5)
            
            title = "BESS" if data['type'] == "BESS" else "SOLAR PV"
            unit = "kWh" if data['type'] == "BESS" else "kWp"
            lines = [title, f"{data['kWac']} kWac", f"{data['cap_val']} {unit}"]
            
            for i, l in enumerate(lines):
                ax.text(px, box_top - 1.2 - (i*0.3), l, color="tab:green", ha="center", fontsize=8, fontweight='bold')
                
        elif itype == 'EMSB':
            by = start_y - 0.8
            ax.plot([px-0.15, px+0.15], [by-0.15, by+0.15], color="tab:blue", lw=2)
            ax.plot([px-0.15, px+0.15], [by+0.15, by-0.15], color="tab:blue", lw=2)
            
            rect = patches.Rectangle((px - 0.6, box_top - 0.6), 1.2, 0.6, fill=True, facecolor="tab:blue", edgecolor="tab:blue")
            ax.add_patch(rect)
            ax.text(px, box_top - 0.9, data, ha="center", va="top", fontsize=8)

def draw_section_feeders_mpl(ax, feeder_indices, x_centers, sub_widths_map, Y_MAIN_BUS, Y_FDR_BRK, swg_configs, swg_names, voltage):
    lv_bus_y = {}; lv_bus_x = {}; lv_bus_edges = {}; sub_feeder_lv_coords = {}
    sub_feeder_bus_edges = {}
    sub_board_bus_coords = {} 
    last_sub_feeder_coords = {} 
    first_sub_feeder_coords = {}
    
    dims = get_feeder_width_config(is_pptx=False)

    for i, cx in zip(feeder_indices, x_centers):
        config = swg_configs.get(i, {})
        ctype = config.get("type", "Standard")
        color = "tab:green" if ctype == "MV Gen" else "tab:blue"
        
        ax.plot([cx, cx], [Y_MAIN_BUS, Y_FDR_BRK], color=color, lw=3)
        ax.plot([cx-0.2, cx+0.2], [Y_FDR_BRK-0.2, Y_FDR_BRK+0.2], color=color, lw=3)
        ax.plot([cx-0.2, cx+0.2], [Y_FDR_BRK+0.2, Y_FDR_BRK-0.2], color=color, lw=3)
        
        if voltage != "400V":
            ax.text(cx, Y_FDR_BRK+0.5, swg_names[i], ha="center", fontsize=10, fontweight="bold")
            
        current_y = Y_FDR_BRK - 0.2
        
        if ctype == "MV Gen":
            gens = config.get("gens", [])
            draw_lv_system_mpl(ax, cx, current_y, gens, False, "")
            
        elif ctype == "Sub-Board":
            sub_voltage = config.get('sub_voltage')
            is_extension = (voltage == sub_voltage)
            
            if not is_extension:
                y_bus_connection = current_y - 1.5
                ax.plot([cx, cx], [current_y, y_bus_connection + 0.6], color="tab:blue", lw=3)
                draw_tx_mpl(ax, cx, y_bus_connection, f"TX-{i+1}", f"{voltage}/{sub_voltage}")
                y_mv_breaker_main = y_bus_connection - 0.7 
                ax.plot([cx, cx], [y_bus_connection - 0.15, y_mv_breaker_main + 0.2], color="tab:blue", lw=3)
                ax.plot([cx-0.15, cx+0.15], [y_mv_breaker_main-0.15, y_mv_breaker_main+0.15], color="tab:blue", lw=3) 
                ax.plot([cx-0.15, cx+0.15], [y_mv_breaker_main+0.15, y_mv_breaker_main-0.15], color="tab:blue", lw=3) 
                y_sub_bus = y_mv_breaker_main - 0.6
                ax.plot([cx, cx], [y_mv_breaker_main - 0.15, y_sub_bus], color="tab:blue", lw=3)
            else:
                y_ext_brk = current_y - 1.5 
                ax.plot([cx, cx], [current_y, y_ext_brk + 0.2], color="tab:blue", lw=3)
                ax.plot([cx-0.2, cx+0.2], [y_ext_brk-0.2, y_ext_brk+0.2], color="tab:blue", lw=3)
                ax.plot([cx-0.2, cx+0.2], [y_ext_brk+0.2, y_ext_brk-0.2], color="tab:blue", lw=3)
                y_sub_bus = y_ext_brk - 0.8
                ax.plot([cx, cx], [y_ext_brk - 0.2, y_sub_bus], color="tab:blue", lw=3)

            sub_feeders = config.get("sub_feeders", {})
            n_subs = len(sub_feeders)
            widths_list = sub_widths_map.get(i, [])
            
            sub_bus_edges_local = {}
            sub_y_local = {}

            if n_subs > 0:
                total_sb_width = sum(widths_list) + (len(widths_list)-1)*dims["sub_gap"]
                start_sub_x = cx - (total_sb_width / 2)
                
                ax.hlines(y_sub_bus, start_sub_x, start_sub_x + total_sb_width, lw=4, color="tab:blue")
                sub_board_bus_coords[i] = (start_sub_x, start_sub_x + total_sb_width, y_sub_bus)
                
                ax.text(start_sub_x + total_sb_width + 0.2, y_sub_bus, sub_voltage, color="tab:blue", fontsize=8, fontweight='bold', va='center', ha='left')
                
                current_sb_x = start_sub_x
                for j in range(n_subs):
                    sub_w = widths_list[j]
                    sub_x = current_sb_x + (sub_w / 2)
                    
                    s_conf = sub_feeders.get(j, {})
                    sf_type = s_conf.get("type", "Standard")
                    s_gens = s_conf.get("gens", [])
                    has_emsb = s_conf.get("has_emsb", False)
                    
                    y_sub_mv_brk = y_sub_bus - 0.8
                    ax.plot([sub_x, sub_x], [y_sub_bus, y_sub_mv_brk+0.15], color="tab:blue", lw=2)
                    ax.plot([sub_x-0.1, sub_x+0.1], [y_sub_mv_brk-0.1, y_sub_mv_brk+0.1], color="tab:blue", lw=2)
                    ax.plot([sub_x-0.1, sub_x+0.1], [y_sub_mv_brk+0.1, y_sub_mv_brk-0.1], color="tab:blue", lw=2)
                    
                    y_end_point = 0
                    x_end_point = sub_x

                    if sf_type == "MV Gen":
                        draw_lv_system_mpl(ax, sub_x, y_sub_mv_brk - 0.2, s_gens, False, "")
                        y_end_point = y_sub_mv_brk - 0.2
                    elif sf_type == "Extension":
                        ext_feeders = s_conf.get("extension_feeders", {})
                        if ext_feeders:
                            # Draw line down to Nested Bus
                            y_nest_bus = y_sub_mv_brk - 1.5
                            ax.plot([sub_x, sub_x], [y_sub_mv_brk-0.15, y_nest_bus], color="tab:blue", lw=2)
                            
                            # Calculate nested width
                            n_ext = len(ext_feeders)
                            total_ext_w = 0
                            ext_item_widths = []
                            for k in ext_feeders:
                                ef_conf = ext_feeders[k]
                                ef_type = ef_conf.get("type", "Standard")
                                if ef_type == "MV Gen":
                                    ef_w = dims["item_w"]
                                else:
                                    ef_gens = ef_conf.get("gens", [])
                                    ef_emsb = ef_conf.get("has_emsb", False)
                                    ef_cnt = len(ef_gens) + (1 if ef_emsb else 0)
                                    ef_w = max(dims["item_w"] * 1.5, ef_cnt * dims["item_w"])
                                ext_item_widths.append(ef_w)
                                total_ext_w += ef_w
                            
                            if n_ext > 1:
                                total_ext_w += (n_ext - 1) * (dims["sub_gap"] * 0.8)
                            
                            nest_start_x = sub_x - total_ext_w/2
                            ax.hlines(y_nest_bus, nest_start_x, nest_start_x + total_ext_w, lw=3, color="tab:blue")
                            
                            # For coupler logic (Extension Bus)
                            sub_bus_edges_local[j] = (nest_start_x, nest_start_x + total_ext_w)
                            sub_y_local[j] = y_nest_bus

                            ax.text(nest_start_x + total_ext_w + 0.1, y_nest_bus, "Ext.", color="tab:blue", fontsize=6)
                            
                            curr_nest_x = nest_start_x
                            
                            ext_first_lv_left = 0
                            ext_first_lv_y = 0
                            ext_last_lv_right = 0
                            ext_last_lv_y = 0
                            
                            nested_lv_coords_mpl = {}

                            for k in range(n_ext):
                                ef_w = ext_item_widths[k]
                                ef_center = curr_nest_x + ef_w/2
                                ef_conf = ext_feeders[k]
                                ef_type = ef_conf.get("type", "Standard")
                                
                                y_nf_brk = y_nest_bus - 0.5
                                ax.plot([ef_center, ef_center], [y_nest_bus, y_nf_brk+0.1], color="tab:blue", lw=1.5)
                                ax.plot([ef_center-0.1, ef_center+0.1], [y_nf_brk-0.1, y_nf_brk+0.1], color="tab:blue", lw=1.5)
                                ax.plot([ef_center-0.1, ef_center+0.1], [y_nf_brk+0.1, y_nf_brk-0.1], color="tab:blue", lw=1.5)
                                
                                this_lv_left = ef_center
                                this_lv_right = ef_center
                                this_lv_y = y_nf_brk 

                                if ef_type == "MV Gen":
                                    draw_lv_system_mpl(ax, ef_center, y_nf_brk - 0.2, ef_conf.get("gens", []), False, "")
                                    this_lv_y = y_nf_brk - 0.2
                                else:
                                    # Standard -> TX
                                    y_nf_tx = y_nf_brk - 0.8
                                    ax.plot([ef_center, ef_center], [y_nf_brk-0.1, y_nf_tx+0.25], color="tab:blue", lw=1.5)
                                    draw_tx_mpl(ax, ef_center, y_nf_tx, "", f"{sub_voltage}/0.4")
                                    
                                    # LV Bus
                                    y_nf_lv = y_nf_tx - 1.0
                                    ax.plot([ef_center, ef_center], [y_nf_tx-0.3, y_nf_lv], color="tab:blue", lw=1.5)
                                    
                                    bus_viz = ef_w - 0.5
                                    ax.hlines(y_nf_lv, ef_center - bus_viz/2, ef_center + bus_viz/2, lw=2, color="tab:blue")
                                    draw_lv_system_mpl(ax, ef_center, y_nf_lv, ef_conf.get("gens", []), ef_conf.get("has_emsb"), "EMSB")
                                    
                                    this_lv_left = ef_center - bus_viz/2
                                    this_lv_right = ef_center + bus_viz/2
                                    this_lv_y = y_nf_lv
                                
                                nested_lv_coords_mpl[k] = (this_lv_left, this_lv_right, this_lv_y)

                                if k == 0:
                                    ext_first_lv_left = this_lv_left
                                    ext_first_lv_y = this_lv_y
                                if k == n_ext - 1:
                                    ext_last_lv_right = this_lv_right
                                    ext_last_lv_y = this_lv_y
                                
                                ax.text(ef_center, y_nf_brk - 0.2, ef_conf.get("name", ""), ha="center", fontsize=6)
                                curr_nest_x += ef_w + (dims["sub_gap"] * 0.8)
                            
                            # Draw internal extension couplers in Preview
                            req_ext_couplers = s_conf.get("extension_couplers", [])
                            for pair_idx in req_ext_couplers:
                                if pair_idx in nested_lv_coords_mpl and (pair_idx+1) in nested_lv_coords_mpl:
                                    r_edge = nested_lv_coords_mpl[pair_idx][1]
                                    y1 = nested_lv_coords_mpl[pair_idx][2]
                                    l_edge = nested_lv_coords_mpl[pair_idx+1][0]
                                    y2 = nested_lv_coords_mpl[pair_idx+1][2]
                                    mid_cy = (y1+y2)/2
                                    
                                    ax.plot([r_edge, l_edge], [mid_cy, mid_cy], color="tab:red", lw=2)
                                    mid_cx = (r_edge + l_edge)/2
                                    ax.plot([mid_cx-0.1, mid_cx+0.1], [mid_cy-0.1, mid_cy+0.1], color="tab:red", lw=2)
                                    ax.plot([mid_cx-0.1, mid_cx+0.1], [mid_cy+0.1, mid_cy-0.1], color="tab:red", lw=2)
                            
                            x_end_point = ext_last_lv_right 
                            y_end_point = ext_last_lv_y
                        else:
                            ax.plot([sub_x, sub_x], [y_sub_mv_brk-0.15, y_sub_mv_brk-1.5], color="tab:blue", lw=2)
                            ax.text(sub_x, y_sub_mv_brk - 1.7, f"{sub_voltage} OUT", ha="center", fontweight="bold")
                            y_end_point = y_sub_mv_brk - 1.5
                            x_end_point = sub_x
                    else:
                        # Standard -> Step Down TX
                        y_tx_sub = y_sub_mv_brk - 1.2
                        ax.plot([sub_x, sub_x], [y_sub_mv_brk-0.15, y_tx_sub+0.3], color="tab:blue", lw=2)
                        draw_tx_mpl(ax, sub_x, y_tx_sub, f"TX-SF{j+1}", f"{sub_voltage}/0.4")
                        y_sub_breaker = y_tx_sub - 1.5
                        ax.plot([sub_x, sub_x], [y_tx_sub - 0.6, y_sub_breaker], color="tab:blue", lw=2)
                        bx, by = sub_x, y_sub_breaker
                        ax.plot([bx-0.1, bx+0.1], [by-0.1, by+0.1], color="tab:blue", lw=2)
                        ax.plot([bx-0.1, bx+0.1], [by+0.1, by-0.1], color="tab:blue", lw=2)
                        y_lv_out = by - 0.2
                        sub_feeder_lv_coords[(i, j)] = (sub_x, y_lv_out)
                        bus_width_viz = max(dims["item_w"], sub_w - 0.5)
                        bus_left = sub_x - bus_width_viz/2
                        bus_right = sub_x + bus_width_viz/2
                        ax.hlines(y_lv_out, bus_left, bus_right, lw=4, color="tab:blue")
                        sub_feeder_bus_edges[(i, j)] = (bus_left, bus_right)
                        sub_bus_edges_local[j] = (bus_left, bus_right)
                        sub_y_local[j] = y_lv_out
                        
                        draw_lv_system_mpl(ax, sub_x, y_lv_out, s_gens, has_emsb, "EMSB")
                        y_end_point = y_lv_out
                        x_end_point = bus_right 
                    
                    if sf_type != "Extension":
                        ax.text(sub_x, y_sub_mv_brk - 0.4 if sf_type=="MV Gen" else y_sub_breaker - 0.4, s_conf.get("name", ""), ha="center", fontsize=7)
                    else:
                        ax.text(sub_x, y_sub_mv_brk - 0.4, s_conf.get("name", ""), ha="center", fontsize=7)
                    
                    if j == 0:
                         if sf_type == "Standard":
                             first_sub_feeder_coords[i] = (sub_x - max(dims["item_w"], sub_w - 0.5)/2, y_end_point)
                         elif sf_type == "Extension" and ext_feeders:
                             first_sub_feeder_coords[i] = (ext_first_lv_left, ext_first_lv_y)

                    if j == n_subs - 1:
                        if sf_type == "Standard":
                             last_sub_feeder_coords[i] = (x_end_point, y_end_point)
                        elif sf_type == "Extension" and ext_feeders:
                             last_sub_feeder_coords[i] = (ext_last_lv_right, ext_last_lv_y)
                    
                    current_sb_x += sub_w + dims["sub_gap"] 

                # Sub Couplers
                for cp in config.get("sub_couplers", []):
                     if cp in sub_bus_edges_local and (cp+1) in sub_bus_edges_local:
                         e1 = sub_bus_edges_local[cp][1]; e2 = sub_bus_edges_local[cp+1][0]; y_cp = sub_y_local[cp]
                         ax.plot([e1, e2], [y_cp, y_cp], color="tab:red", lw=2)
                         mid_cx = (e1+e2)/2
                         ax.plot([mid_cx-0.1, mid_cx+0.1], [y_cp-0.1, y_cp+0.1], color="tab:red", lw=2)
                         ax.plot([mid_cx-0.1, mid_cx+0.1], [y_cp+0.1, y_cp-0.1], color="tab:red", lw=2)

        else: # Standard
            chain = get_tx_chain(voltage, config.get("tx_scheme", ""))
            temp_y = current_y
            if not chain and voltage == "400V":
                y_bus = 2.0; ax.plot([cx, cx], [temp_y, y_bus], color="tab:blue", lw=3)
                lv_bus_y[i] = y_bus; lv_bus_x[i] = cx
                y_fin = y_bus
            else:
                for step in chain:
                    y_tx = temp_y - 1.5
                    draw_tx_mpl(ax, cx, y_tx, f"TX-{i+1}", step["ratio"])
                    ax.plot([cx, cx], [temp_y, y_tx+0.6], color="tab:blue", lw=3, zorder=1)
                    temp_y = y_tx - 0.6
                y_bus = temp_y - 1.0
                ax.plot([cx, cx], [temp_y, y_bus], color="tab:blue", lw=3)
                lv_bus_y[i] = y_bus; lv_bus_x[i] = cx
                y_fin = y_bus
            
            gens = config.get("gens", []); has_emsb = config.get("emsb", {}).get("has")
            cnt = len(gens) + (1 if has_emsb else 0)
            bw = max(dims["min_w"], cnt * dims["item_w"])
            
            left_edge = cx - int(bw/2); right_edge = cx + int(bw/2)
            ax.hlines(lv_bus_y[i], left_edge, right_edge, lw=5, color="tab:blue")
            lv_bus_edges[i] = (left_edge, right_edge)
            ax.text(cx, lv_bus_y[i]-0.3, config.get("msb_name", ""), ha="center", va="top", fontweight="bold")
            draw_lv_system_mpl(ax, cx, lv_bus_y[i], gens, has_emsb, config["emsb"]["name"])

            # Store endpoints for Inter-Feeder Coupling (Standard Feeders)
            first_sub_feeder_coords[i] = (left_edge, lv_bus_y[i])
            last_sub_feeder_coords[i] = (right_edge, lv_bus_y[i])

    return lv_bus_x, lv_bus_y, lv_bus_edges, sub_feeder_lv_coords, sub_feeder_bus_edges, sub_board_bus_coords, last_sub_feeder_coords, first_sub_feeder_coords

def draw_feeder_group_on_slide(slide, voltage, feeders_list, swg_configs, swg_names, 
                               start_x, dims, incomer_data, 
                               draw_bc_start, draw_bc_end, bc_label, 
                               scale_factor):
    """
    Draws a group of feeders. All dimensions are raw inches * scale_factor.
    Returns dictionaries where coordinates are int (EMU).
    """
    
    Y_MAIN_BUS = int(S(6.0, scale_factor))
    Y_INC_TOP = int(S(1.0, scale_factor))
    Y_INC_BRK = int(S(4.0, scale_factor))
    Y_FDR_BRK = int(S(7.5, scale_factor))
    GAP = int(S(dims["gap"], scale_factor))
    
    current_x = int(start_x)
    
    feeder_widths = []
    total_group_width = 0
    
    for idx in feeders_list:
        conf = swg_configs.get(idx, {})
        w_raw, sub_ws_raw = calculate_single_feeder_width(conf, dims)
        w_scaled = int(S(w_raw, scale_factor))
        sub_ws_scaled = [int(S(sw, scale_factor)) for sw in sub_ws_raw]
        
        feeder_widths.append((w_scaled, sub_ws_scaled))
        total_group_width += w_scaled + GAP
        
    bus_left = int(start_x)
    
    if draw_bc_start:
        add_continuation_arrow(slide, bus_left, Y_MAIN_BUS, "prev", "From Sheet 1", scale_factor)
        bus_left += int(S(0.8, scale_factor))
        add_line(slide, bus_left, Y_MAIN_BUS, bus_left + int(S(0.8, scale_factor)), Y_MAIN_BUS, 3, RGBColor(255,0,0))
        bus_left += int(S(0.8, scale_factor))
        
    actual_bus_end = bus_left + total_group_width - GAP
    
    if draw_bc_end:
        add_busbar(slide, bus_left, Y_MAIN_BUS, actual_bus_end - bus_left)
        c_start = actual_bus_end
        ext_len = int(S(2.5, scale_factor))
        c_end = c_start + ext_len
        add_line(slide, c_start, Y_MAIN_BUS, c_end, Y_MAIN_BUS, 3, RGBColor(255,0,0))
        mid_bc = (c_start + c_end) // 2
        add_breaker_x(slide, mid_bc, Y_MAIN_BUS, scale_factor, 0.25, RGBColor(255,0,0))
        tb = slide.shapes.add_textbox(int(mid_bc - S(1.5, scale_factor)), int(Y_MAIN_BUS - S(1.5, scale_factor)), int(S(3.0, scale_factor)), int(S(1.2, scale_factor)))
        tb.text_frame.text = bc_label
        for p in tb.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(255,0,0)
            p.font.size = Pt(max(10, 20*scale_factor)) 
            p.alignment = PP_ALIGN.CENTER
        add_continuation_arrow(slide, c_end, Y_MAIN_BUS, "next", "To Sheet 2", scale_factor)
    else:
        add_busbar(slide, bus_left, Y_MAIN_BUS, actual_bus_end - bus_left)

    cursor_x = bus_left 
    
    if incomer_data:
        inc_x = bus_left + (actual_bus_end - bus_left) // 2
        add_line(slide, inc_x, Y_INC_TOP, inc_x, Y_MAIN_BUS)
        add_breaker_x(slide, inc_x, Y_INC_BRK, scale_factor, 0.3)
        tb_w = int(S(6.0, scale_factor))
        tb_x = inc_x - tb_w // 2
        tb_y = int(Y_INC_TOP - S(1.0, scale_factor))
        tb = slide.shapes.add_textbox(tb_x, tb_y, tb_w, int(S(1.5, scale_factor)))
        tb.text_frame.text = incomer_data["label"]
        for p in tb.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(max(10, 20*scale_factor)) 

    lv_coords_local = {}
    sub_board_bus_local = {}
    last_sub_local = {}
    first_sub_local = {}

    for i, idx in enumerate(feeders_list):
        config = swg_configs.get(idx, {})
        w_feeder_scaled, sub_ws_scaled = feeder_widths[i]
        
        cx = cursor_x + w_feeder_scaled // 2
        
        ctype = config.get("type", "Standard")
        col = RGBColor(0,176,80) if ctype == "MV Gen" else RGBColor(0,112,192)

        add_line(slide, cx, Y_MAIN_BUS, cx, int(Y_FDR_BRK + S(0.1, scale_factor)), 3, col)
        add_breaker_x(slide, cx, Y_FDR_BRK, scale_factor, 0.25, col)
        
        if voltage != "400V":
            tb = slide.shapes.add_textbox(int(cx - S(2, scale_factor)), int(Y_FDR_BRK - S(0.8, scale_factor)), int(S(4, scale_factor)), int(S(0.8, scale_factor)))
            tb.text_frame.text = swg_names[idx]
            p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.bold = True; p.font.size = Pt(max(8, 16*scale_factor)) 

        cur_y = int(Y_FDR_BRK + S(0.1, scale_factor))
        y_fin_this_feeder = 0; lv_edges = (0, 0)

        if ctype == "MV Gen":
            gens = config.get("gens", [])
            pptx_add_inverter_branch(slide, cx, cur_y, gens, scale_factor)
            
        elif ctype == "Sub-Board":
            sub_voltage = config.get('sub_voltage')
            is_extension = (voltage == sub_voltage)
            
            y_tx1 = int(cur_y + S(2.5, scale_factor))
            y_sub_bus = int(y_tx1 + S(2.7, scale_factor))
            
            if not is_extension:
                pptx_add_transformer(slide, cx, y_tx1, f"{voltage}/{sub_voltage}", f"TX-{idx+1}", scale_factor)
                add_line(slide, cx, cur_y, cx, int(y_tx1 - S(0.9, scale_factor)))
                
                y_mv_breaker_main = int(y_tx1 + S(1.5, scale_factor))
                add_line(slide, cx, int(y_tx1 + S(0.9, scale_factor)), cx, int(y_mv_breaker_main - S(0.2, scale_factor)))
                add_breaker_x(slide, cx, y_mv_breaker_main, scale_factor, 0.25)
                
                y_sub_bus = int(y_mv_breaker_main + S(1.2, scale_factor))
                add_line(slide, cx, int(y_mv_breaker_main + S(0.2, scale_factor)), cx, int(y_sub_bus + S(0.05, scale_factor)))
            else:
                y_ext_breaker = y_tx1
                add_line(slide, cx, cur_y, cx, int(y_ext_breaker - S(0.2, scale_factor)))
                add_breaker_x(slide, cx, y_ext_breaker, scale_factor, 0.25)
                y_sub_bus = int(y_ext_breaker + S(1.2, scale_factor))
                add_line(slide, cx, int(y_ext_breaker + S(0.2, scale_factor)), cx, int(y_sub_bus + S(0.05, scale_factor)))
                
            sub_feeders = config.get("sub_feeders", {})
            n_subs = len(sub_feeders)
            if n_subs > 0:
                total_sb_width = sum(sub_ws_scaled) + (len(sub_ws_scaled)-1)*int(S(dims["sub_gap"], scale_factor))
                start_sub_x = cx - total_sb_width // 2
                add_busbar(slide, start_sub_x, y_sub_bus, total_sb_width)
                
                # Storing int (EMU)
                sub_board_bus_local[idx] = (start_sub_x, start_sub_x + total_sb_width, y_sub_bus)
                
                tb = slide.shapes.add_textbox(start_sub_x + total_sb_width, int(y_sub_bus - S(0.3, scale_factor)), int(S(1.0, scale_factor)), int(S(0.5, scale_factor)))
                tb.text_frame.text = sub_voltage; tb.text_frame.paragraphs[0].font.size = Pt(max(10, 20*scale_factor))

                curr_sb_x = start_sub_x
                sub_bus_edges_local = {}; sub_y_local = {}

                for j in range(n_subs):
                    sw = sub_ws_scaled[j]; sx = curr_sb_x + sw // 2
                    s_conf = sub_feeders.get(j, {})
                    sf_type = s_conf.get("type", "Standard")
                    
                    y_mv_brk_sub = int(y_sub_bus + S(1.2, scale_factor))
                    add_line(slide, sx, y_sub_bus, sx, int(y_mv_brk_sub - S(0.2, scale_factor)))
                    add_breaker_x(slide, sx, y_mv_brk_sub, scale_factor, 0.2)
                    
                    y_end_pt = 0 
                    x_end_pt = 0
                    
                    if sf_type == "MV Gen":
                        pptx_add_inverter_branch(slide, sx, int(y_mv_brk_sub + S(0.2, scale_factor)), s_conf.get("gens", []), scale_factor)
                        y_end_pt = int(y_mv_brk_sub + S(2.0, scale_factor))
                        x_end_pt = sx 
                    elif sf_type == "Extension":
                        ext_feeders = s_conf.get("extension_feeders", {})
                        if ext_feeders:
                            y_nest_bus = int(y_mv_brk_sub + S(2.5, scale_factor))
                            add_line(slide, sx, int(y_mv_brk_sub + S(0.2, scale_factor)), sx, y_nest_bus)
                            
                            total_ext_w = 0
                            ext_item_widths = []
                            for k in ext_feeders:
                                ef_conf = ext_feeders[k]
                                ef_type = ef_conf.get("type", "Standard")
                                if ef_type == "MV Gen":
                                    ef_w = int(S(dims["item_w"], scale_factor))
                                else:
                                    ef_gens = ef_conf.get("gens", [])
                                    ef_emsb = ef_conf.get("has_emsb", False)
                                    ef_cnt = len(ef_gens) + (1 if ef_emsb else 0)
                                    ef_w = max(int(S(dims["item_w"] * 1.5, scale_factor)), ef_cnt * int(S(dims["item_w"], scale_factor)))
                                ext_item_widths.append(ef_w)
                                total_ext_w += ef_w
                            
                            if len(ext_feeders) > 1:
                                total_ext_w += (len(ext_feeders) - 1) * int(S(dims["sub_gap"] * 0.8, scale_factor))
                                
                            nest_start_x = sx - total_ext_w // 2
                            add_busbar(slide, nest_start_x, y_nest_bus, total_ext_w)
                            
                            sub_bus_edges_local[j] = (nest_start_x, nest_start_x + total_ext_w)
                            sub_y_local[j] = y_nest_bus

                            curr_nest_x = nest_start_x
                            
                            ext_first_lv_left = 0
                            ext_first_lv_y = 0
                            ext_last_lv_right = 0
                            ext_last_lv_y = 0
                            
                            nested_lv_coords = {} 

                            n_ext = len(ext_feeders)
                            for k in range(n_ext):
                                ef_conf = ext_feeders[k]
                                ef_type = ef_conf.get("type", "Standard")
                                ef_w = ext_item_widths[k]
                                ef_center = curr_nest_x + ef_w // 2
                                
                                y_nf_brk = int(y_nest_bus + S(1.0, scale_factor))
                                add_line(slide, ef_center, y_nest_bus, ef_center, y_nf_brk)
                                add_breaker_x(slide, ef_center, y_nf_brk, scale_factor, 0.15)
                                
                                this_lv_left = ef_center
                                this_lv_right = ef_center
                                this_lv_y = y_nf_brk 
                                
                                if ef_type == "MV Gen":
                                    pptx_add_inverter_branch(slide, ef_center, int(y_nf_brk + S(0.1, scale_factor)), ef_conf.get("gens", []), scale_factor)
                                    this_lv_y = int(y_nf_brk + S(2.0, scale_factor))
                                else:
                                    y_nf_tx = int(y_nf_brk + S(1.5, scale_factor))
                                    add_line(slide, ef_center, int(y_nf_brk + S(0.1, scale_factor)), ef_center, int(y_nf_tx - S(0.9, scale_factor)))
                                    pptx_add_transformer(slide, ef_center, y_nf_tx, f"{sub_voltage}/0.4", "", scale_factor)
                                    
                                    y_nf_lv = int(y_nf_tx + S(1.5, scale_factor))
                                    add_line(slide, ef_center, int(y_nf_tx + S(0.9, scale_factor)), ef_center, y_nf_lv)
                                    
                                    bus_viz = ef_w - int(S(0.5, scale_factor))
                                    add_busbar(slide, ef_center - bus_viz // 2, y_nf_lv, bus_viz)
                                    pptx_add_lv_system(slide, ef_center, y_nf_lv, ef_conf.get("gens", []), ef_conf.get("has_emsb"), "EMSB", scale_factor)
                                    
                                    this_lv_left = ef_center - bus_viz // 2
                                    this_lv_right = ef_center + bus_viz // 2
                                    this_lv_y = y_nf_lv
                                
                                nested_lv_coords[k] = (this_lv_left, this_lv_right, this_lv_y)

                                if k == 0:
                                    ext_first_lv_left = this_lv_left
                                    ext_first_lv_y = this_lv_y
                                if k == n_ext - 1:
                                    ext_last_lv_right = this_lv_right
                                    ext_last_lv_y = this_lv_y
                                
                                tb = slide.shapes.add_textbox(int(ef_center - S(1.0, scale_factor)), int(y_nf_brk - S(0.5, scale_factor)), int(S(2.0, scale_factor)), int(S(0.5, scale_factor)))
                                tb.text_frame.text = ef_conf.get("name", "")
                                tb.text_frame.paragraphs[0].font.size = Pt(max(8, 10*scale_factor))
                                tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                                
                                curr_nest_x += ef_w + int(S(dims["sub_gap"] * 0.8, scale_factor))
                            
                            # Draw internal extension couplers (PPTX)
                            req_ext_couplers = s_conf.get("extension_couplers", [])
                            for pair_idx in req_ext_couplers:
                                if pair_idx in nested_lv_coords and (pair_idx+1) in nested_lv_coords:
                                    r_edge = nested_lv_coords[pair_idx][1]
                                    y1 = nested_lv_coords[pair_idx][2]
                                    l_edge = nested_lv_coords[pair_idx+1][0]
                                    y2 = nested_lv_coords[pair_idx+1][2]
                                    
                                    mid_cy = (y1 + y2) // 2
                                    
                                    add_line(slide, r_edge, mid_cy, l_edge, mid_cy, 3, RGBColor(255,0,0))
                                    mid_cx = (r_edge + l_edge) // 2
                                    add_breaker_x(slide, mid_cx, mid_cy, scale_factor, 0.15, RGBColor(255,0,0))

                            y_end_pt = ext_last_lv_y 
                            x_end_pt = ext_last_lv_right 
                        else:
                            add_line(slide, sx, int(y_mv_brk_sub + S(0.2, scale_factor)), sx, int(y_mv_brk_sub + S(2.5, scale_factor)))
                            tb = slide.shapes.add_textbox(int(sx - S(1.5, scale_factor)), int(y_mv_brk_sub + S(2.7, scale_factor)), int(S(3.0, scale_factor)), int(S(0.8, scale_factor)))
                            tb.text_frame.text = f"{sub_voltage} OUT"
                            tb.text_frame.paragraphs[0].font.size = Pt(max(10, 14*scale_factor))
                            tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            tb.text_frame.paragraphs[0].font.bold = True
                            y_end_pt = int(y_mv_brk_sub + S(3.0, scale_factor))
                            x_end_pt = sx
                    else:
                        y_tx_sub = int(y_mv_brk_sub + S(2.2, scale_factor))
                        add_line(slide, sx, int(y_mv_brk_sub + S(0.2, scale_factor)), sx, int(y_tx_sub - S(0.9, scale_factor)))
                        pptx_add_transformer(slide, sx, y_tx_sub, f"{sub_voltage}/0.4", f"TX-SF{j+1}", scale_factor)
                        
                        y_sub_breaker = int(y_tx_sub + S(1.5, scale_factor))
                        add_line(slide, sx, int(y_tx_sub + S(0.9, scale_factor)), sx, y_sub_breaker)
                        add_breaker_x(slide, sx, y_sub_breaker, scale_factor, 0.2)
                        
                        y_lv_out = int(y_sub_breaker + S(0.2, scale_factor))
                        b_viz = max(int(S(dims["item_w"], scale_factor)), sw - int(S(0.5, scale_factor)))
                        b_start = sx - b_viz // 2
                        b_end = b_start + b_viz
                        add_busbar(slide, b_start, y_lv_out, b_viz)
                        
                        if j == n_subs - 1 and i == len(feeders_list) - 1:
                            tb = slide.shapes.add_textbox(int(b_end + S(0.1, scale_factor)), int(y_lv_out - S(0.3, scale_factor)), int(S(1.5, scale_factor)), int(S(0.6, scale_factor)))
                            tb.text_frame.text = "400V"
                            p = tb.text_frame.paragraphs[0]
                            p.font.size = Pt(max(10, 20*scale_factor)) # FIXED 20pt
                            p.font.bold = True
                            p.font.color.rgb = RGBColor(0, 112, 192)
                            p.alignment = PP_ALIGN.LEFT
                        
                        sub_bus_edges_local[j] = (b_start, b_start + b_viz)
                        sub_y_local[j] = y_lv_out
                        
                        pptx_add_lv_system(slide, sx, y_lv_out, s_conf.get("gens", []), s_conf.get("has_emsb"), "EMSB", scale_factor)
                        
                        y_end_pt = y_lv_out
                        x_end_pt = b_end
                    
                    curr_sb_x += sw + int(S(dims["sub_gap"], scale_factor))
                    
                    if j == 0:
                         if sf_type == "Standard":
                             first_sub_local[idx] = (int(sx - max(S(dims["item_w"], scale_factor), sw - S(0.5, scale_factor))/2), y_end_pt)
                         elif sf_type == "Extension" and ext_feeders:
                             first_sub_local[idx] = (ext_first_lv_left, ext_first_lv_y)

                    if j == n_subs - 1:
                        if sf_type == "Standard":
                             last_sub_local[idx] = (x_end_pt, y_end_pt)
                        elif sf_type == "Extension" and ext_feeders:
                             last_sub_local[idx] = (ext_last_lv_right, ext_last_lv_y)
                
                for cp in config.get("sub_couplers", []):
                     if cp in sub_bus_edges_local and (cp+1) in sub_bus_edges_local:
                         e1 = sub_bus_edges_local[cp][1]; e2 = sub_bus_edges_local[cp+1][0]; y_cp = sub_y_local[cp]
                         add_line(slide, e1, y_cp, e2, y_cp, 3, RGBColor(255,0,0))
                         add_breaker_x(slide, (e1+e2) // 2, y_cp, scale_factor, 0.2, RGBColor(255,0,0))

        else: # Standard
            chain = get_tx_chain(voltage, config.get("tx_scheme", ""))
            temp_y = cur_y
            if not chain and voltage == "400V":
                y_fin_this_feeder = int(S(14.0, scale_factor))
                add_line(slide, cx, temp_y, cx, y_fin_this_feeder)
            else:
                for step in chain:
                    y_tx = int(temp_y + S(2.5, scale_factor))
                    pptx_add_transformer(slide, cx, y_tx, step["ratio"], f"TX-{idx+1}", scale_factor)
                    add_line(slide, cx, temp_y, cx, int(y_tx - S(0.9, scale_factor)))
                    temp_y = int(y_tx + S(0.9, scale_factor))
                y_fin_this_feeder = int(temp_y + S(2.0, scale_factor))
                add_line(slide, cx, temp_y, cx, int(y_fin_this_feeder + S(0.05, scale_factor)))
            
            gens = config.get("gens", []); has_emsb = config.get("emsb", {}).get("has")
            cnt = len(gens) + (1 if has_emsb else 0)
            bw = max(int(S(dims["min_w"], scale_factor)), cnt * int(S(dims["item_w"], scale_factor)))
            
            left_edge = cx - bw // 2; right_edge = cx + bw // 2
            add_busbar(slide, left_edge, y_fin_this_feeder, bw)
            
            if i == len(feeders_list) - 1:
                tb = slide.shapes.add_textbox(int(right_edge + S(0.1, scale_factor)), int(y_fin_this_feeder - S(0.3, scale_factor)), int(S(1.5, scale_factor)), int(S(0.6, scale_factor)))
                tb.text_frame.text = "400V"
                p = tb.text_frame.paragraphs[0]
                p.font.size = Pt(max(10, 20*scale_factor)) 
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 112, 192)
                p.alignment = PP_ALIGN.LEFT
            
            lv_edges = (left_edge, right_edge)
            pptx_add_lv_system(slide, cx, y_fin_this_feeder, gens, has_emsb, config["emsb"]["name"], scale_factor)
            
            first_sub_local[idx] = (left_edge, y_fin_this_feeder)
            last_sub_local[idx] = (right_edge, y_fin_this_feeder)

        if ctype == "Standard":
            lv_coords_local[idx] = {"y": y_fin_this_feeder, "left": lv_edges[0], "right": lv_edges[1]}

        cursor_x += w_feeder_scaled + GAP

    drawn_width = cursor_x - start_x
    return lv_coords_local, drawn_width, actual_bus_end, sub_board_bus_local, last_sub_local, first_sub_local


def generate_pptx(voltage, num_in, num_swg, section_distribution, inc_bc_status, 
                  msb_bc_status, lv_couplers, lv_bc_status, swg_names, swg_configs,
                  inter_sub_bus_couplers=None, inter_lv_couplers=None):
    
    if inter_sub_bus_couplers is None: inter_sub_bus_couplers = []
    if inter_lv_couplers is None: inter_lv_couplers = []

    prs = Presentation()
    MAX_PPTX_WIDTH_INCHES = 56.0 
    dims = get_feeder_width_config(is_pptx=True)
    GAP_RAW = dims["gap"]
    
    sections = []
    curr = 0
    for count in section_distribution:
        if count > 0:
            sections.append(list(range(curr, curr + count)))
        curr += count
    
    section_raw_widths = []
    has_sub_board = False
    for sec_indices in sections:
        total_w = 0
        for idx in sec_indices:
            conf = swg_configs.get(idx, {})
            if conf.get("type") == "Sub-Board": has_sub_board = True
            w, _ = calculate_single_feeder_width(conf, dims)
            total_w += w + GAP_RAW
        section_raw_widths.append(total_w)
    
    total_raw_width = sum(section_raw_widths) + (len(sections)-1)*2.0 
    needed_width = total_raw_width + 4.0
    requires_split = needed_width > MAX_PPTX_WIDTH_INCHES

    global_lv_map = {}
    global_sub_bus_map = {}
    global_last_sub = {}
    global_first_sub = {}
    
    needed_height = 30.0 if has_sub_board else 20.0
    Y_MAIN_BUS_CONST = 6.0 

    # --- SCENARIO A: SINGLE SLIDE ---
    if not requires_split:
        scale = 1.0
        final_w_inches = max(20.0, needed_width)
        
        prs.slide_width = int(Inches(final_w_inches))
        prs.slide_height = int(Inches(needed_height))
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        start_margin = int(Inches(final_w_inches - total_raw_width) / 2)
        current_x = start_margin
        
        for s_i, feeders in enumerate(sections):
            inc_label = f"INCOMING {s_i+1}\n({voltage})"
            is_last = (s_i == len(sections) - 1)
            
            coords, w_used, bus_end_x, sb_loc, l_loc, f_loc = draw_feeder_group_on_slide(slide, voltage, feeders, swg_configs, swg_names, 
                                                current_x, dims, 
                                                {"label": inc_label}, 
                                                False, False, "", scale)
            
            for idx, data in coords.items(): 
                data['slide'] = slide
                data['scale'] = scale 
                global_lv_map[idx] = data
            
            for k, v in sb_loc.items(): global_sub_bus_map[k] = {'coords': v, 'slide': slide, 'scale': scale}
            for k, v in l_loc.items(): global_last_sub[k] = {'coords': v, 'slide': slide, 'scale': scale}
            for k, v in f_loc.items(): global_first_sub[k] = {'coords': v, 'slide': slide, 'scale': scale}

            current_x += w_used
            
            if not is_last:
                gap_size = int(S(1.0, scale))
                y_bus = int(S(Y_MAIN_BUS_CONST, scale))
                add_line(slide, current_x - int(S(GAP_RAW, scale)), y_bus, current_x + gap_size, y_bus, 3, RGBColor(255,0,0))
                mid_x = (current_x + gap_size // 2 - int(S(GAP_RAW, scale)) // 2)
                add_breaker_x(slide, mid_x, y_bus, scale, 0.25, RGBColor(255,0,0))
                bc_text = f"BC-{s_i+1}\n({msb_bc_status.get(s_i, 'NO')})"
                tb_bc = slide.shapes.add_textbox(int(mid_x - S(1.5, scale)), int(y_bus - S(1.2, scale)), int(S(3.0, scale)), int(S(0.8, scale)))
                tb_bc.text_frame.text = bc_text
                for p in tb_bc.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.color.rgb = RGBColor(255,0,0)
                    p.font.size = Pt(max(10, 20*scale)) 

                current_x += gap_size
            else:
                tb = slide.shapes.add_textbox(int(bus_end_x + S(0.1, scale)), int(S(Y_MAIN_BUS_CONST, scale) - S(0.3, scale)), int(S(1.5, scale)), int(S(0.6, scale)))
                tb.text_frame.text = voltage
                p = tb.text_frame.paragraphs[0]
                p.font.size = Pt(max(10, 20*scale)) 
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 112, 192)
                p.alignment = PP_ALIGN.LEFT

    # --- SCENARIO B: SPLIT ---
    else:
        lhs_raw_w = 0
        rhs_raw_w = 0
        lhs_fds = []
        rhs_fds = []
        
        if len(sections) == 1:
            all_fds = sections[0]
            mid = len(all_fds)//2
            lhs_fds = all_fds[:mid]
            rhs_fds = all_fds[mid:]
            
            for idx in lhs_fds: lhs_raw_w += calculate_single_feeder_width(swg_configs[idx], dims)[0] + GAP_RAW
            lhs_raw_w += 2.0 
            
            for idx in rhs_fds: rhs_raw_w += calculate_single_feeder_width(swg_configs[idx], dims)[0] + GAP_RAW
            rhs_raw_w += 2.0
            
        else:
            lhs_fds = sections[0]
            lhs_raw_w = section_raw_widths[0] + 2.0
            rhs_indices_groups = sections[1:]
            rhs_raw_w = sum(section_raw_widths[1:]) + (len(rhs_indices_groups)-1)*1.0 + 2.0

        max_content_w = max(lhs_raw_w, rhs_raw_w)
        target_slide_w = max_content_w + 2.0
        final_slide_w = min(target_slide_w, MAX_PPTX_WIDTH_INCHES)
        final_slide_w = max(final_slide_w, 20.0) 
        
        prs.slide_width = int(Inches(final_slide_w))
        prs.slide_height = int(Inches(needed_height))
        
        available_w = final_slide_w - 2.0
        scale_lhs_w = min(1.0, available_w / lhs_raw_w)
        scale_rhs_w = min(1.0, available_w / rhs_raw_w)
        scale_h_limit = 0.85
        
        scale_lhs = min(scale_lhs_w, scale_h_limit)
        scale_rhs = min(scale_rhs_w, scale_h_limit)
        
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        lhs_content_w = S(lhs_raw_w, scale_lhs)
        start_x1 = int((Inches(final_slide_w) - lhs_content_w) / 2)
        
        if len(sections) == 1:
             c1, _, _, sb1, l1, f1 = draw_feeder_group_on_slide(slide1, voltage, lhs_fds, swg_configs, swg_names, 
                                            start_x1, dims, {"label": f"INCOMING 1\n({voltage})"}, 
                                            False, True, "Bus Cont.", scale_lhs)
        else:
             c1, _, _, sb1, l1, f1 = draw_feeder_group_on_slide(slide1, voltage, lhs_fds, swg_configs, swg_names, 
                                            start_x1, dims, {"label": f"INCOMING 1\n({voltage})"}, 
                                            False, True, f"BC-1\n({msb_bc_status.get(0, 'NO')})", scale_lhs)
        for idx, data in c1.items(): 
            data['slide'] = slide1
            data['scale'] = scale_lhs
            global_lv_map[idx] = data
        for k, v in sb1.items(): global_sub_bus_map[k] = {'coords': v, 'slide': slide1, 'scale': scale_lhs}
        for k, v in l1.items(): global_last_sub[k] = {'coords': v, 'slide': slide1, 'scale': scale_lhs}
        for k, v in f1.items(): global_first_sub[k] = {'coords': v, 'slide': slide1, 'scale': scale_lhs}
        
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        rhs_content_w = S(rhs_raw_w, scale_rhs)
        start_x2 = int((Inches(final_slide_w) - rhs_content_w) / 2)
        
        if len(sections) == 1:
             c2, _, bus_end_x2, sb2, l2, f2 = draw_feeder_group_on_slide(slide2, voltage, rhs_fds, swg_configs, swg_names, 
                                            start_x2, dims, {"label": ""}, 
                                            True, False, "", scale_rhs)
             for idx, data in c2.items(): 
                 data['slide'] = slide2
                 data['scale'] = scale_rhs
                 global_lv_map[idx] = data
             for k, v in sb2.items(): global_sub_bus_map[k] = {'coords': v, 'slide': slide2, 'scale': scale_rhs}
             for k, v in l2.items(): global_last_sub[k] = {'coords': v, 'slide': slide2, 'scale': scale_rhs}
             for k, v in f2.items(): global_first_sub[k] = {'coords': v, 'slide': slide2, 'scale': scale_rhs}
             
             tb = slide2.shapes.add_textbox(int(bus_end_x2 + S(0.1, scale_rhs)), int(S(Y_MAIN_BUS_CONST, scale_rhs) - S(0.3, scale_rhs)), int(S(1.5, scale_rhs)), int(S(0.6, scale_rhs)))
             tb.text_frame.text = voltage
             p = tb.text_frame.paragraphs[0]
             p.font.size = Pt(max(10, 20*scale_rhs)) 
             p.font.bold = True
             p.font.color.rgb = RGBColor(0, 112, 192)
             p.alignment = PP_ALIGN.LEFT

        else:
             curr_x = start_x2
             rhs_indices_groups = sections[1:]
             for r_i, r_feeders in enumerate(rhs_indices_groups):
                is_first = (r_i == 0)
                real_inc_idx = r_i + 2 
                lbl = f"INCOMING {real_inc_idx}\n({voltage})"
                
                c_out, w_used, bus_end_x, sb_out, l_out, f_out = draw_feeder_group_on_slide(slide2, voltage, r_feeders, swg_configs, swg_names, 
                                                    curr_x, dims, {"label": lbl},
                                                    is_first, False, "", scale_rhs)
                for idx, data in c_out.items(): 
                    data['slide'] = slide2
                    data['scale'] = scale_rhs
                    global_lv_map[idx] = data
                for k, v in sb_out.items(): global_sub_bus_map[k] = {'coords': v, 'slide': slide2, 'scale': scale_rhs}
                for k, v in l_out.items(): global_last_sub[k] = {'coords': v, 'slide': slide2, 'scale': scale_rhs}
                for k, v in f_out.items(): global_first_sub[k] = {'coords': v, 'slide': slide2, 'scale': scale_rhs}
                
                if r_i < len(rhs_indices_groups) - 1:
                    curr_x += w_used
                    add_line(slide2, curr_x - int(S(GAP_RAW, scale_rhs)), int(S(6.0, scale_rhs)), curr_x + int(S(1.0, scale_rhs)), int(S(6.0, scale_rhs)), 3, RGBColor(255,0,0))
                    mid_x = (curr_x + int(S(1.0, scale_rhs)) // 2 - int(S(GAP_RAW, scale_rhs)) // 2)
                    add_breaker_x(slide2, mid_x, int(S(6.0, scale_rhs)), scale_rhs, 0.25, RGBColor(255,0,0))
                    bc_text = f"BC-{r_i+2}\n({msb_bc_status.get(r_i+1, 'NO')})"
                    tb_bc = slide2.shapes.add_textbox(int(mid_x - S(0.5, scale_rhs)), int(S(6.0, scale_rhs) - S(1.2, scale_rhs)), int(S(3.0, scale_rhs)), int(S(0.8, scale_rhs)))
                    tb_bc.text_frame.text = bc_text
                    for p in tb_bc.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        p.font.color.rgb = RGBColor(255,0,0)
                        p.font.size = Pt(max(10, 20*scale_rhs)) 

                    curr_x += int(S(1.0, scale_rhs))
                else:
                    tb = slide2.shapes.add_textbox(int(bus_end_x + S(0.1, scale_rhs)), int(S(Y_MAIN_BUS_CONST, scale_rhs) - S(0.3, scale_rhs)), int(S(1.5, scale_rhs)), int(S(0.6, scale_rhs)))
                    tb.text_frame.text = voltage
                    p = tb.text_frame.paragraphs[0]
                    p.font.size = Pt(max(10, 20*scale_rhs)) 
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 112, 192)
                    p.alignment = PP_ALIGN.LEFT

    # 5. 11kV Sub-Board Couplers
    for pair_idx in inter_sub_bus_couplers:
        f1 = pair_idx; f2 = pair_idx + 1
        if f1 in global_sub_bus_map and f2 in global_sub_bus_map:
             d1 = global_sub_bus_map[f1]; d2 = global_sub_bus_map[f2]
             
             x1 = d1['coords'][1] # right edge of left bus
             x2 = d2['coords'][0] # left edge of right bus
             y1 = d1['coords'][2]
             y2 = d2['coords'][2]
             mid_y = (y1 + y2) // 2
             sc = d1['scale']

             if d1['slide'] == d2['slide']:
                 sl = d1['slide']
                 if abs(y1 - y2) > 1000:
                    add_line(sl, x1, y1, x1, mid_y, 3, RGBColor(255,0,0))
                    add_line(sl, x1, mid_y, x2, mid_y, 3, RGBColor(255,0,0))
                    add_line(sl, x2, mid_y, x2, y2, 3, RGBColor(255,0,0))
                 else:
                    add_line(sl, x1, mid_y, x2, mid_y, 3, RGBColor(255,0,0)) 
                 
                 mid_x = (x1 + x2) // 2
                 add_breaker_x(sl, mid_x, mid_y, sc, 0.25, RGBColor(255, 0, 0))

                 # Text
                 tb = sl.shapes.add_textbox(int(mid_x - S(1.0, sc)), int(mid_y + S(0.2, sc)), int(S(2.0, sc)), int(S(0.8, sc)))
                 tb.text_frame.text = "BC (11kV)"
                 tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
                 tb.text_frame.paragraphs[0].font.size = Pt(max(10, 14*sc))
                 tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
             else:
                 # Split case
                 sl1 = d1['slide']; sl2 = d2['slide']
                 
                 # Slide 1: Line -> Breaker -> Arrow
                 ext_len = int(S(2.0, sc))
                 end_x1 = x1 + ext_len
                 add_line(sl1, x1, y1, end_x1, y1, 3, RGBColor(255,0,0))
                 mid_x1 = (x1 + end_x1) // 2
                 add_breaker_x(sl1, mid_x1, y1, sc, 0.25, RGBColor(255,0,0))
                 add_continuation_arrow(sl1, end_x1, y1, "next", "To Next Bus", sc)
                 
                 tb1 = sl1.shapes.add_textbox(int(mid_x1 - S(1.0, sc)), int(y1 + S(0.2, sc)), int(S(2.0, sc)), int(S(0.8, sc)))
                 tb1.text_frame.text = "BC"
                 tb1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
                 tb1.text_frame.paragraphs[0].font.size = Pt(max(10, 14*sc))
                 tb1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                 # Slide 2: Arrow -> Line
                 sc2 = d2['scale']
                 start_x2 = x2 - int(S(1.5, sc2))
                 add_continuation_arrow(sl2, start_x2, y2, "prev", "From Prev Bus", sc2)
                 add_line(sl2, start_x2, y2, x2, y2, 3, RGBColor(255,0,0))

    # 6. 0.4kV Inter-Feeder Couplers
    for pair_idx in inter_lv_couplers:
        f1 = pair_idx; f2 = pair_idx + 1
        if f1 in global_last_sub and f2 in global_first_sub:
             d1 = global_last_sub[f1]; d2 = global_first_sub[f2]
             
             x1 = d1['coords'][0]
             x2 = d2['coords'][0]
             y1 = d1['coords'][1]
             y2 = d2['coords'][1]
             mid_y = (y1 + y2) // 2
             sc = d1['scale']
             
             if d1['slide'] == d2['slide']:
                 sl = d1['slide']
                 
                 if abs(y1 - y2) > 1000:
                    add_line(sl, x1, y1, x1, mid_y, 3, RGBColor(255,0,0))
                    add_line(sl, x1, mid_y, x2, mid_y, 3, RGBColor(255,0,0))
                    add_line(sl, x2, mid_y, x2, y2, 3, RGBColor(255,0,0))
                 else:
                    add_line(sl, x1, mid_y, x2, mid_y, 3, RGBColor(255,0,0)) 
                 
                 mid_x = (x1 + x2) // 2
                 add_breaker_x(sl, mid_x, mid_y, sc, 0.2, RGBColor(255, 0, 0))

                 tb = sl.shapes.add_textbox(int(mid_x - S(1.0, sc)), int(mid_y + S(0.2, sc)), int(S(2.0, sc)), int(S(0.8, sc)))
                 tb.text_frame.text = "LV-BC"
                 tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
                 tb.text_frame.paragraphs[0].font.size = Pt(max(10, 14*sc))
                 tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
             else:
                 # Split case
                 sl1 = d1['slide']; sl2 = d2['slide']
                 
                 # Slide 1: Line -> Breaker -> Arrow
                 ext_len = int(S(2.0, sc))
                 end_x1 = x1 + ext_len
                 add_line(sl1, x1, y1, end_x1, y1, 3, RGBColor(255,0,0))
                 mid_x1 = (x1 + end_x1) // 2
                 add_breaker_x(sl1, mid_x1, y1, sc, 0.2, RGBColor(255,0,0))
                 add_continuation_arrow(sl1, end_x1, y1, "next", "To Next LV", sc)
                 
                 tb1 = sl1.shapes.add_textbox(int(mid_x1 - S(1.0, sc)), int(y1 + S(0.2, sc)), int(S(2.0, sc)), int(S(0.8, sc)))
                 tb1.text_frame.text = "LV-BC"
                 tb1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
                 tb1.text_frame.paragraphs[0].font.size = Pt(max(10, 14*sc))
                 tb1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                 # Slide 2: Arrow -> Line
                 sc2 = d2['scale']
                 start_x2 = x2 - int(S(1.5, sc2))
                 add_continuation_arrow(sl2, start_x2, y2, "prev", "From Prev LV", sc2)
                 add_line(sl2, start_x2, y2, x2, y2, 3, RGBColor(255,0,0))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

def draw_preview_mpl(voltage, num_in, num_swg, section_distribution, inc_bc_status, 
                     msb_bc_status, lv_couplers, lv_bc_status, swg_names, swg_configs, 
                     inter_sub_bus_couplers=None, inter_lv_couplers=None):
    
    if inter_sub_bus_couplers is None: inter_sub_bus_couplers = []
    if inter_lv_couplers is None: inter_lv_couplers = []

    start_margin = 2.0
    total_w = 0
    temp_x = start_margin
    # Calculate total width first to size figure
    for count in section_distribution:
        if count == 0: continue 
        w, _, _, _ = calculate_section_layout(range(count), swg_configs, temp_x, is_pptx=False)
        total_w += w + 3.0
        temp_x += w + 3.0
        
    fig_w = max(12.0, total_w + 5.0) 
    fig, ax = plt.subplots(figsize=(fig_w, 16))
    
    Y_INC_TOP = 13; Y_INC_BRK = 12; Y_MAIN_BUS = 11; Y_FDR_BRK = 10
    
    current_feeder_idx = 0
    current_x_start = start_margin
    
    global_lv_x = {}; global_lv_y = {}; global_lv_edges = {}; global_sub_coords = {}
    global_sub_bus_edges = {}
    global_sub_board_bus_coords = {}
    global_last_sub_coords = {}
    global_first_sub_coords = {}

    bus_endpoints = []
    
    for s_idx, count in enumerate(section_distribution):
        if count == 0: 
            bus_endpoints.append((None, None)); continue

        indices = list(range(current_feeder_idx, current_feeder_idx + count))
        width, centers, _, sub_w_map = calculate_section_layout(indices, swg_configs, current_x_start, is_pptx=False)
        
        bus_min = min(centers) - 3.0; bus_max = max(centers) + 3.0
        ax.hlines(Y_MAIN_BUS, bus_min, bus_max, lw=6, color="tab:blue")
        bus_endpoints.append((bus_min, bus_max))
        
        inc_x = (bus_min + bus_max) / 2 if num_in == 1 else (bus_max - 1.5 if s_idx == 0 else bus_min + 1.5)
        ax.plot([inc_x, inc_x], [Y_INC_TOP, Y_MAIN_BUS], color="tab:blue", lw=3)
        ax.plot([inc_x-0.2, inc_x+0.2], [Y_INC_BRK-0.2, Y_INC_BRK+0.2], color="tab:blue", lw=3)
        ax.plot([inc_x-0.2, inc_x+0.2], [Y_INC_BRK+0.2, Y_INC_BRK-0.2], color="tab:blue", lw=3)
        ax.text(inc_x, Y_INC_TOP + 0.2, f"INCOMING {s_idx+1}\n({voltage})", ha="center", fontweight="bold")
        
        if s_idx == len(section_distribution) - 1:
             ax.text(bus_max + 0.5, Y_MAIN_BUS, voltage, ha="left", va="center", fontweight="bold", fontsize=12, color="tab:blue")
        
        lv_x, lv_y, lv_edges, sub_c, sub_edges, sb_bus_coords, last_sb_c, first_sb_c = draw_section_feeders_mpl(ax, indices, centers, sub_w_map, Y_MAIN_BUS, Y_FDR_BRK, swg_configs, swg_names, voltage)
        
        global_lv_x.update(lv_x); global_lv_y.update(lv_y); global_lv_edges.update(lv_edges)
        global_sub_coords.update(sub_c); global_sub_bus_edges.update(sub_edges)
        global_sub_board_bus_coords.update(sb_bus_coords)
        global_last_sub_coords.update(last_sb_c)
        global_first_sub_coords.update(first_sb_c)
        
        current_feeder_idx += count
        current_x_start += width + 3.0
        
    # Draw Main Bus Couplers
    for s_idx in range(len(section_distribution) - 1):
        left_sect = bus_endpoints[s_idx]
        right_sect = bus_endpoints[s_idx+1]
        if left_sect[1] is not None and right_sect[0] is not None:
            ax.plot([left_sect[1], right_sect[0]], [Y_MAIN_BUS, Y_MAIN_BUS], color="tab:red", lw=3)
            mid_c = (left_sect[1] + right_sect[0]) / 2
            ax.plot([mid_c-0.2, mid_c+0.2], [Y_MAIN_BUS-0.2, Y_MAIN_BUS+0.2], color="tab:red", lw=3)
            ax.plot([mid_c-0.2, mid_c+0.2], [Y_MAIN_BUS+0.2, Y_MAIN_BUS-0.2], color="tab:red", lw=3)
            ax.text(mid_c, Y_MAIN_BUS + 0.5, f"BC-{s_idx+1}\n({msb_bc_status.get(s_idx, 'NO')})", color="tab:red", ha="center", fontweight="bold")

    # Draw Inter-Feeder Bus Couplers (11kV / 33kV)
    for idx in inter_sub_bus_couplers:
        if idx in global_sub_board_bus_coords and (idx+1) in global_sub_board_bus_coords:
            _, right_edge_1, y1 = global_sub_board_bus_coords[idx]
            left_edge_2, _, y2 = global_sub_board_bus_coords[idx+1]
            mid_y = (y1+y2)/2
            
            # Draw vertical segments if heights differ
            ax.plot([right_edge_1, right_edge_1], [y1, mid_y], color="tab:red", lw=3)
            ax.plot([right_edge_1, left_edge_2], [mid_y, mid_y], color="tab:red", lw=3)
            ax.plot([left_edge_2, left_edge_2], [mid_y, y2], color="tab:red", lw=3)
            
            mid_x = (right_edge_1 + left_edge_2)/2
            ax.plot([mid_x-0.2, mid_x+0.2], [mid_y-0.2, mid_y+0.2], color="tab:red", lw=3)
            ax.plot([mid_x-0.2, mid_x+0.2], [mid_y+0.2, mid_y-0.2], color="tab:red", lw=3)
            ax.text(mid_x, mid_y+0.5, "BC (11kV)", color="tab:red", ha="center", fontsize=8, fontweight="bold")

    # Draw Inter-Feeder LV Couplers (0.4kV)
    for idx in inter_lv_couplers:
        if idx in global_last_sub_coords and (idx+1) in global_first_sub_coords:
            x1, y1 = global_last_sub_coords[idx]
            x2, y2 = global_first_sub_coords[idx+1]
            mid_y = (y1+y2)/2
            
            ax.plot([x1, x1], [y1, mid_y], color="tab:red", lw=3)
            ax.plot([x1, x2], [mid_y, mid_y], color="tab:red", lw=3)
            ax.plot([x2, x2], [mid_y, y2], color="tab:red", lw=3)
            
            mid_x = (x1+x2)/2
            ax.plot([mid_x-0.2, mid_x+0.2], [mid_y-0.2, mid_y+0.2], color="tab:red", lw=3)
            ax.plot([mid_x-0.2, mid_x+0.2], [mid_y+0.2, mid_y-0.2], color="tab:red", lw=3)
            ax.text(mid_x, mid_y+0.5, "LV-BC", color="tab:red", ha="center", fontsize=8, fontweight="bold")

    ax.axis('off'); ax.set_ylim(-8, 14); ax.set_xlim(0, fig_w)
    return fig

def main():
    st.set_page_config(layout="wide", page_title="SLD Generator")
    
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False

    if not st.session_state.authenticated:
        st.title("🔒 Internal Access Only")
        passcode = st.text_input("Enter Passcode:", type="password")
        if st.button("Login"):
            if passcode == "9999":
                st.session_state.authenticated = True
                st.rerun()
            else:
                st.error("Incorrect Passcode")
        return 

    st.title("⚡ SLD Generator")

    with st.sidebar:
        st.subheader("System Configuration")
        if st.button("Reset All", type="secondary"):
            for key in list(st.session_state.keys()):
                if key != "authenticated":
                    del st.session_state[key]
            st.rerun()

        voltage = st.selectbox("Voltage", ["400V", "11kV", "33kV", "132kV"], key="sys_v")
        num_in = st.selectbox("Incomers", [1, 2, 3], index=1, key="sys_in")
        n_swg = st.number_input("Total Number of Feeders", 1, 20, 4)
        
        section_distribution = []
        if num_in == 1:
            section_distribution = [n_swg]
        else:
            st.markdown("### Bus Section Configuration")
            remaining = n_swg
            for i in range(num_in - 1):
                val = st.number_input(f"Feeders on Bus Section {i+1}", 0, remaining, max(1, remaining//2), key=f"sec_{i}")
                section_distribution.append(val); remaining -= val
            section_distribution.append(remaining)
            st.info(f"Feeders on Bus Section {num_in}: {remaining}")

        msb_bc_status = {}
        inc_bc_status = [] # Initialize safety
        
        if num_in > 1:
            st.markdown("### Bus Coupler Status")
            for i in range(num_in - 1):
                msb_bc_status[i] = st.selectbox(f"Bus Coupler {i+1}-{i+2}", ["NO", "NC"], key=f"mbc_{i}")

        swg_configs = {}; swg_names = []
        
        st.markdown("### Feeder Details")
        for i in range(n_swg):
            with st.expander(f"Feeder {i+1}", expanded=False):
                name = st.text_input("Name", f"F-{i+1}", key=f"n_{i}")
                swg_names.append(name)
                
                valid_types = ["Standard", "MV Gen"]
                if voltage not in ["400V", "11kV"]: valid_types.append("Sub-Board")
                if voltage != "400V": valid_types.append("Extension")
                
                ctype = "Standard"
                if voltage != "400V": ctype = st.selectbox("Type", valid_types, key=f"t_{i}")
                
                gens = []; has_emsb = False
                
                if ctype == "MV Gen":
                    gens, has_emsb = get_mv_gen_inputs(f"g_{i}"), False
                elif ctype == "Standard":
                    gens, has_emsb = get_lv_gen_inputs(f"g_{i}", include_emsb=True)
                
                conf = {"type": ctype, "msb_name": name, "gens": gens, "emsb": {"has": has_emsb, "name": "EMSB"}}
                
                if ctype == "Standard":
                    conf["tx_scheme"] = f"{voltage}/0.4 kV"
                    
                elif ctype == "Sub-Board":
                    conf["sub_voltage"] = st.selectbox("Sub Voltage", ["11kV", "6.6kV"], key=f"sv_{i}")
                    n_sub_feeders = st.number_input(f"No. of {conf['sub_voltage']} Feeders", 1, 10, 2, key=f"nsf_{i}")
                    
                    if n_sub_feeders > 1:
                        valid_s_pairs = list(range(n_sub_feeders - 1))
                        s_labels = [f"SF-{p+1} & SF-{p+2}" for p in valid_s_pairs]
                        sel_s_couplers = st.multiselect("Add LV Coupler between:", s_labels, key=f"ssc_{i}")
                        conf["sub_couplers"] = [valid_s_pairs[s_labels.index(l)] for l in sel_s_couplers]
                    
                    conf["sub_feeders"] = {}
                    for j in range(n_sub_feeders):
                        st.caption(f"Sub-Feeder {j+1}")
                        sf_name = st.text_input(f"Name", f"SF-{j+1}", key=f"sfn_{i}_{j}")
                        
                        sf_type_options = ["Standard", "MV Gen", "Extension"]
                        sf_type = st.selectbox("Sub-Feeder Type", sf_type_options, key=f"sft_{i}_{j}")
                        
                        sf_gens = []
                        sf_emsb = False
                        
                        if sf_type == "Standard":
                            sf_gens, sf_emsb = get_lv_gen_inputs(f"sfg_{i}_{j}", include_emsb=True)
                        elif sf_type == "MV Gen":
                            sf_gens, sf_emsb = get_mv_gen_inputs(f"sfg_{i}_{j}"), False
                        
                        ext_feeders_data = {}
                        ext_couplers = [] # Initialize here
                        if sf_type == "Extension":
                            n_ext = st.number_input(f"No. of Feeders on Ext {j+1}", 1, 10, 2, key=f"next_{i}_{j}")
                            
                            if n_ext > 1:
                                valid_ie_pairs = list(range(n_ext - 1))
                                ie_labels = [f"EF-{p+1} & EF-{p+2}" for p in valid_ie_pairs]
                                sel_ie_couplers = st.multiselect("Add Coupler Inside Extension:", ie_labels, key=f"ie_c_{i}_{j}")
                                ext_couplers = [valid_ie_pairs[ie_labels.index(l)] for l in sel_ie_couplers]

                            for k in range(n_ext):
                                st.markdown(f"**Ext Feeder {k+1}**")
                                ef_name = st.text_input(f"Name", f"EF-{k+1}", key=f"efn_{i}_{j}_{k}")
                                ef_type = st.selectbox("Type", ["Standard", "MV Gen"], key=f"eft_{i}_{j}_{k}")
                                ef_gens = []
                                ef_emsb = False
                                if ef_type == "Standard":
                                    ef_gens, ef_emsb = get_lv_gen_inputs(f"efg_{i}_{j}_{k}", True)
                                else:
                                    ef_gens, ef_emsb = get_mv_gen_inputs(f"efg_{i}_{j}_{k}"), False
                                ext_feeders_data[k] = {"type": ef_type, "name": ef_name, "gens": ef_gens, "has_emsb": ef_emsb}

                        conf["sub_feeders"][j] = {"type": sf_type, "name": sf_name, "gens": sf_gens, "has_emsb": sf_emsb, "extension_feeders": ext_feeders_data, "extension_couplers": ext_couplers}
                
                elif ctype == "Extension":
                    conf["type"] = "Sub-Board" 
                    conf["sub_voltage"] = voltage
                    
                    st.info(f"Extension at {voltage}")
                    n_sub_feeders = st.number_input(f"No. of Feeders on Extension", 1, 10, 2, key=f"nef_{i}")
                    
                    if n_sub_feeders > 1:
                        valid_e_pairs = list(range(n_sub_feeders - 1))
                        e_labels = [f"EF-{p+1} & EF-{p+2}" for p in valid_e_pairs]
                        sel_e_couplers = st.multiselect("Add Bus Coupler between:", e_labels, key=f"ext_bc_{i}")
                        conf["sub_couplers"] = [valid_e_pairs[e_labels.index(l)] for l in sel_e_couplers]
                    
                    conf["sub_feeders"] = {}
                    for j in range(n_sub_feeders):
                        st.caption(f"Extension Feeder {j+1}")
                        sf_name = st.text_input(f"Name", f"EF-{j+1}", key=f"efn_{i}_{j}")
                        
                        sf_type = st.selectbox("Type", ["Standard", "MV Gen"], key=f"eft_{i}_{j}")
                        
                        sf_gens = []
                        sf_emsb = False
                        
                        if sf_type == "Standard":
                             sf_gens, sf_emsb = get_lv_gen_inputs(f"efg_{i}_{j}", include_emsb=True)
                        else:
                             sf_gens, sf_emsb = get_mv_gen_inputs(f"efg_{i}_{j}"), False
                             
                        conf["sub_feeders"][j] = {"type": sf_type, "name": sf_name, "gens": sf_gens, "has_emsb": sf_emsb}

                swg_configs[i] = conf

        lv_couplers = []; lv_bc_status = {}
        with st.expander("LV (0.4kV) Bus Couplers"):
            valid_pairs = []
            for i in range(n_swg-1):
                if swg_configs[i]["type"] == "Standard" and swg_configs[i+1]["type"] == "Standard":
                    valid_pairs.append(i)
            if valid_pairs:
                pair_lbls = [f"#{p+1} & #{p+2}" for p in valid_pairs]
                sel_lv = st.multiselect("Couples", pair_lbls, key="lv_c_sel")
                for s in sel_lv:
                    idx = pair_lbls.index(s); real_idx = valid_pairs[idx]
                    lv_couplers.append(real_idx); lv_bc_status[real_idx] = st.selectbox(f"Status {s}", ["NO", "NC"], key=f"lvbc_{real_idx}")
            else:
                st.write("No adjacent standard feeders available for coupling.")

        inter_sub_bus_couplers = []
        inter_lv_couplers = []
        
        if voltage != "400V":
            with st.expander("Inter-Feeder Bus Couplers (11kV & 0.4kV)"):
                sb_pairs = []
                for i in range(n_swg - 1):
                    c1 = swg_configs[i]; c2 = swg_configs[i+1]
                    if c1["type"] == "Sub-Board" and c2["type"] == "Sub-Board":
                        if c1.get("sub_voltage") == c2.get("sub_voltage"):
                             sb_pairs.append(i)
                
                if sb_pairs:
                    sb_labels = [f"F-{p+1} & F-{p+2} ({swg_configs[p].get('sub_voltage')})" for p in sb_pairs]
                    sel_sb = st.multiselect("Select Intermediate Bus Couplers", sb_labels, key="inter_sb_c")
                    for s in sel_sb:
                        idx = sb_labels.index(s)
                        inter_sub_bus_couplers.append(sb_pairs[idx])
                
                ilv_pairs = []
                for i in range(n_swg - 1):
                     c1 = swg_configs[i]; c2 = swg_configs[i+1]
                     has_lv_1 = False; has_lv_2 = False
                     
                     if c1["type"] == "Standard": has_lv_1 = True
                     elif c1["type"] == "Sub-Board":
                         subs = c1.get("sub_feeders", {})
                         if subs:
                             last_idx = len(subs) - 1
                             if subs[last_idx].get("type") == "Standard": has_lv_1 = True
                             elif subs[last_idx].get("type") == "Extension": has_lv_1 = True
                     
                     if c2["type"] == "Standard": has_lv_2 = True
                     elif c2["type"] == "Sub-Board":
                         subs = c2.get("sub_feeders", {})
                         if subs:
                             if subs[0].get("type") == "Standard": has_lv_2 = True
                             elif subs[0].get("type") == "Extension": has_lv_2 = True
                     
                     if has_lv_1 and has_lv_2:
                         ilv_pairs.append(i)
                         
                if ilv_pairs:
                    ilv_labels = [f"F-{p+1} & F-{p+2} (0.4kV)" for p in ilv_pairs]
                    sel_ilv = st.multiselect("Select Inter-Feeder LV Couplers", ilv_labels, key="inter_lv_c")
                    for s in sel_ilv:
                        idx = ilv_labels.index(s)
                        inter_lv_couplers.append(ilv_pairs[idx])


    st.subheader("Preview")
    fig = draw_preview_mpl(voltage, num_in, n_swg, section_distribution, inc_bc_status, 
                           msb_bc_status, lv_couplers, lv_bc_status, swg_names, swg_configs, 
                           inter_sub_bus_couplers, inter_lv_couplers)
    st.pyplot(fig, use_container_width=True)
    
    plt.close(fig)
    
    pptx_data = generate_pptx(voltage, num_in, n_swg, section_distribution, inc_bc_status, msb_bc_status, 
                              lv_couplers, lv_bc_status, swg_names, swg_configs,
                              inter_sub_bus_couplers, inter_lv_couplers)
    
    st.download_button("📥 Download PowerPoint", pptx_data, 
                       f"SLD_{voltage}.pptx", 
                       "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                       type="primary", use_container_width=True)

if __name__ == "__main__":
    main()
